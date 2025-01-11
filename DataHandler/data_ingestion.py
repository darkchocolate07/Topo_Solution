import os
import pandas as pd
import pdfplumber
from pptx import Presentation


class DataIngestion:
    def __init__(self):
        self.data = {}

    def _check_extension(self, filepath, expected_extension):
        """
        Helper method to validate the file extension.
        """
        _, extension = os.path.splitext(filepath)
        if extension.lower() != expected_extension:
            raise ValueError(
                f"Invalid file extension for {filepath}. Expected '{expected_extension}', but got '{extension}'."
            )
        
    def _check_file_exists(self, filepath):
        """
        Checks if the file exists.
        """
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"File {filepath} does not exist.")

    def load_json(self, filepath):
        """
        Load and process a JSON file.
        """
        self._check_file_exists(filepath)
        self._check_extension(filepath, ".json")
        try:
            dataset = pd.read_json(filepath)
            json_data = dataset.iloc[:, 0]
            self.data["json"] = pd.concat(
                [
                    pd.json_normalize(row["employees"]).assign(
                        company_id=row["id"],
                        company_name=row["name"],
                        revenue=row["revenue"]
                    )
                    for row in json_data
                ],
                ignore_index=True,
            )
        except Exception as e:
            raise RuntimeError(f"Error loading JSON file: {e}")

    def load_csv(self, filepath):
        """
        Load and process a CSV file.
        """
        self._check_file_exists(filepath)
        self._check_extension(filepath, ".csv")
        try:
            self.data["csv"] = pd.read_csv(filepath)
        except Exception as e:
            raise RuntimeError(f"Error loading CSV file: {e}")

    def load_pdf(self, filepath):
        """
        Load and process a PDF file.
        """
        self._check_file_exists(filepath)
        self._check_extension(filepath, ".pdf")
        try:
            with pdfplumber.open(filepath) as pdf:
                first_page = pdf.pages[0]
                table = first_page.extract_table()
            columns, rows = table[0], table[1:]
            self.data["pdf"] = pd.DataFrame(rows, columns=columns)
        except Exception as e:
            raise RuntimeError(f"Error loading PDF file: {e}")

    def load_ppt(self, filepath):
        """
        Load and process a PowerPoint file.
        """
        self._check_file_exists(filepath)
        self._check_extension(filepath, ".pptx")
        try:
            presentation = Presentation(filepath)
            slide = presentation.slides[1]  # Slide 2
            table_data = [
                [cell.text for cell in row.cells]
                for shape in slide.shapes if shape.has_table
                for row in shape.table.rows
            ]
            columns, rows = table_data[0], table_data[1:]
            self.data["ppt"] = pd.DataFrame(rows, columns=columns)
        except Exception as e:
            raise RuntimeError(f"Error loading PPTX file: {e}")

    def get_data(self):
        """
        Get the loaded data.
        """
        return self.data
